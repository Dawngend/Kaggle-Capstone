# Copyright 2026 Google LLC
#
# Licensed under the Apache License, Version 2.0 (the "License");
# you may not use this file except in compliance with the License.
# You may obtain a copy of the License at
#
#     https://www.apache.org/licenses/LICENSE-2.0
#
# Unless required by applicable law or agreed to in writing, software
# distributed under the License is distributed on an "AS IS" BASIS,
# WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
# See the License for the specific language governing permissions and
# limitations under the License.

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9 ratio
    prs.slide_height = Inches(7.5)

    # Styling colors
    DARK_BLUE = RGBColor(15, 23, 42)    # Slate 900
    LIGHT_CREAM = RGBColor(250, 246, 240)  # Cream background
    ACCENT_BLUE = RGBColor(37, 99, 235)  # Blue accent
    TEXT_DARK = RGBColor(15, 23, 42)
    TEXT_MUTED = RGBColor(71, 85, 105)

    # Slide 1: Title Slide
    slide_layout = prs.slide_layouts[6] # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Set background color
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = LIGHT_CREAM

    # Add Title
    title_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.333), Inches(2.0))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Runtime Terrors Companion"
    p.font.name = "Outfit"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    p.alignment = PP_ALIGN.LEFT

    # Add Subtitle
    p2 = tf.add_paragraph()
    p2.text = "Multi-Agent AI Study & Athletic Performance Companion"
    p2.font.name = "Inter"
    p2.font.size = Pt(24)
    p2.font.color.rgb = TEXT_DARK
    p2.font.italic = True
    p2.alignment = PP_ALIGN.LEFT

    # Add Metadata (Authors & Track)
    meta_box = slide.shapes.add_textbox(Inches(1.0), Inches(4.5), Inches(11.333), Inches(1.5))
    tf_meta = meta_box.text_frame
    tf_meta.word_wrap = True
    
    p3 = tf_meta.paragraphs[0]
    p3.text = "Kaggle Capstone Project: Agents for Good Track"
    p3.font.name = "Inter"
    p3.font.size = Pt(16)
    p3.font.bold = True
    p3.font.color.rgb = TEXT_MUTED

    p4 = tf_meta.add_paragraph()
    p4.text = "Author: Dawn Andrei Pamesa (BS Data Science student at FEU Institute of Technology)"
    p4.font.name = "Inter"
    p4.font.size = Pt(14)
    p4.font.color.rgb = TEXT_MUTED

    # Set Slide Notes (Script)
    slide.notes_slide.notes_text_frame.text = (
        "Presenter Script:\n"
        "\"For student-athletes, coaching assistants, and high-performance performers, "
        "balancing intensive athletic training with rigorous data science coursework is a severe friction point. "
        "When we spar or run drills, physical exhaustion directly impacts our cognitive bandwidth, leading to burnout. "
        "To solve this, we built the 'Runtime Terrors Companion'—a multi-agent AI system designed to dynamically balance physical recovery with structured academic learning pipelines. Our mission is to keep student-athletes physically healthy while helping them master complex data science concepts.\""
    )

    # Helper function to generate standardized slides
    def add_content_slide(title, points, script_text):
        new_slide = prs.slides.add_slide(prs.slide_layouts[6])
        new_slide.background.fill.solid()
        new_slide.background.fill.fore_color.rgb = LIGHT_CREAM

        # Title Box
        t_box = new_slide.shapes.add_textbox(Inches(1.0), Inches(0.8), Inches(11.333), Inches(1.0))
        tf_title = t_box.text_frame
        tf_title.word_wrap = True
        p_t = tf_title.paragraphs[0]
        p_t.text = title
        p_t.font.name = "Outfit"
        p_t.font.size = Pt(40)
        p_t.font.bold = True
        p_t.font.color.rgb = ACCENT_BLUE

        # Content Box
        c_box = new_slide.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.333), Inches(4.5))
        tf_content = c_box.text_frame
        tf_content.word_wrap = True

        for i, pt in enumerate(points):
            if i == 0:
                p_c = tf_content.paragraphs[0]
            else:
                p_c = tf_content.add_paragraph()
            p_c.text = "• " + pt
            p_c.font.name = "Inter"
            p_c.font.size = Pt(18)
            p_c.font.color.rgb = TEXT_DARK
            p_c.space_after = Pt(14)
        
        new_slide.notes_slide.notes_text_frame.text = f"Presenter Script:\n{script_text}"
        return new_slide

    # Slide 2: The Friction
    add_content_slide(
        title="The Friction: Physical Load vs. Cognitive Fatigue",
        points=[
            "Friction Point: Balancing intense athletic workouts with demanding BS Data Science courses.",
            "Burnout Impact: High training loads (sparring, drills) lead to physical fatigue and mental fog.",
            "Solution: A system that dynamically adapts study loads to match the user's current athletic recovery state.",
            "Track: Agents for Good - optimizing academic learning performance while safeguarding physical student health."
        ],
        script_text=(
            "\"Let's talk about the core friction. Student-athletes have to manage rigorous coursework alongside heavy "
            "physical training. High training volume leads to mental burnout and fatigue. The Runtime Terrors Companion "
            "solves this by tracking training load and adjusting study schedules dynamically. This falls perfectly "
            "into the 'Agents for Good' track, safeguarding student health while supporting academic outcomes.\""
        )
    )

    # Slide 3: Multi-Agent Architecture
    add_content_slide(
        title="ADK Multi-Agent Architecture",
        points=[
            "Academic Tutor Agent: Elite FEU Tech DS tutor breaking down complex ML algorithms & checking comprehension.",
            "Performance Coach Agent: Coordinates athletic sessions, computes recovery metrics, and guides load balance.",
            "Orchestration Hub: Custom orchestrator managing conversation contexts via relational SQLite/PostgreSQL databases.",
            "Interoperable Protocol: Exposes core components to external clients via Model Context Protocol (MCP) Server."
        ],
        script_text=(
            "\"Our system architecture is centered around a multi-agent orchestrator built with the Agent Development Kit. "
            "The Academic Tutor Agent guides data science students, while the Performance Coach Agent logs training metrics "
            "and recovery data. These agents maintain user context through a relational database adapter. Finally, the "
            "entire system is exposed via a Model Context Protocol server, letting external clients execute our tools.\""
        )
    )

    # Slide 4: Key Course Concepts Demonstrated
    add_content_slide(
        title="Key Course Concepts Demonstrated",
        points=[
            "Agent / Multi-Agent System (ADK): Coordinates Tutor and Coach agents with fallback layers in orchestrator.py.",
            "MCP Server: Exposes recovery index calculations and study note parser to external AI tools in mcp_server.py.",
            "Agent Skills: Standing mathematical, document parsing (PDF/PPTX), and 5-Phase ML Sandbox in skills.py.",
            "Security: Credentials managed securely; database queries fully parameterized to prevent SQL injection.",
            "Deployability: Packaged via multi-architecture Dockerfile and configured for PostgreSQL serverless runtimes."
        ],
        script_text=(
            "\"This project serves as a showcase of the key concepts from the course. We've implemented a robust "
            "multi-agent hub using the ADK, a FastMCP server that exposes our core skills, standalone OCR and mathematical "
            "skills, security layers to exclude raw API keys and prevent SQL injection, and a production-ready Docker "
            "configuration for cloud deployment.\""
        )
    )

    # Slide 5: Core Specialized Skills
    add_content_slide(
        title="Core Specialized Agent Skills",
        points=[
            "Sophy Engine: Dynamic spaced repetition card generator (Taglish/Tagalog support) utilizing SuperMemo SM-2 algorithm.",
            "Document Parser & OCR: Converts uploaded PDF/PPTX slides to clean text; fallbacks to Tesseract OCR for scanned modules.",
            "5-Phase Machine Learning Sandbox: Zero-leakage EDA, ColumnTransformer prep, CV, and XGBoost/LightGBM classification.",
            "Recovery Load Calculator: Calculates training load (minutes x RPE) and maps recovery recommendations."
        ],
        script_text=(
            "\"Under the hood, the agents call standalone mathematical and algorithmic skills. Sophy generates taglish "
            "flashcards scheduled with the SuperMemo SM-2 algorithm based on recall accuracy. Our document parser extracts "
            "text from PDF and PPTX files with OCR fallbacks. And our ML Sandbox lets students run a 5-phase classification "
            "pipeline securely without any data leakage.\""
        )
    )

    # Slide 6: Production Deployability & Security
    add_content_slide(
        title="Production Security & Deployability",
        points=[
            "GCP Application Credentials: Authenticates Vertex AI clients securely without hardcoded API keys.",
            "Parameterized Queries: Prevents SQL injection across SQLite and PostgreSQL databases via DatabaseAdapter.",
            "Docker Container: Packages Python runtime, LightGBM, and Tesseract OCR dependencies for Cloud Run execution.",
            "Leakage-free Pipelines: Restricts preprocessing fitting within Stratified K-Fold splits to preserve test sets."
        ],
        script_text=(
            "\"In production environments, security is non-negotiable. The companion avoids hardcoded API keys by utilizing "
            "Application Default Credentials. Database operations are strictly bound using parameterized queries to block SQL "
            "injection. The system is containerized using Docker, making it highly deployable to serverless container hosts "
            "like Google Cloud Run. Thank you!\""
        )
    )

    output_path = os.path.join("submission", "presentation.pptx")
    prs.save(output_path)
    print(f"Presentation saved successfully to: {output_path}")

if __name__ == "__main__":
    create_deck()
