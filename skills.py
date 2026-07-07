# Copyright 2026 Google LLC
#
# Licensed under the Apache License, Version 2.0 (the "License");
# you may not use this file except in compliance with the License.
# You may obtain a copy of the License at
#
#     https://www.apache.org/licenses/LICENSE-2.0
#
# Unless required by applicable law or agreed to in writing, software
# distributed under the License is distributed on an "AS IS" BASIS,
# WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
# See the License for the specific language governing permissions and
# limitations under the License.

import re
import os
import datetime
from typing import Dict, Any, List, Tuple
import pandas as pd
import numpy as np

# Machine Learning Library Imports
from sklearn.compose import ColumnTransformer
from sklearn.pipeline import Pipeline
from sklearn.impute import SimpleImputer
from sklearn.preprocessing import StandardScaler, OneHotEncoder
from sklearn.model_selection import StratifiedKFold, cross_validate
from sklearn.metrics import classification_report, confusion_matrix
from xgboost import XGBClassifier
from lightgbm import LGBMClassifier

# Document Parsing Imports
import pdfplumber
import pptx
import pytesseract
from PIL import Image

# ----------------- COGNITIVE ATHLETIC & ACADEMIC SKILLS -----------------

def calculate_recovery_index(duration_minutes: int, intensity_rpe: int) -> Dict[str, Any]:
    """Calculates training load and computes a dynamic recovery index (0 to 100)."""
    load = duration_minutes * intensity_rpe
    penalty = 0.0
    if intensity_rpe >= 8:
        penalty = (intensity_rpe - 7) * 5.0
        
    raw_recovery = 100.0 - (load / 12.0) - penalty
    recovery_score = max(0.0, min(100.0, raw_recovery))
    
    if recovery_score < 40.0:
        recommendation = "CRITICAL LOAD: Complete rest recommended. Skip high-focus academic or intensive training tasks."
    elif recovery_score < 70.0:
        recommendation = "MODERATE LOAD: Recommended study session limited to 1.5 hours. Perform low-impact drills."
    else:
        recommendation = "OPTIMAL RECOVERY: Full capacity. Ready for core machine learning exercises and intensive sparring."
        
    return {
        "training_load": load,
        "recovery_score": round(recovery_score, 2),
        "recommendation": recommendation
    }

def parse_study_markdown(content: str) -> List[Dict[str, Any]]:
    """Parses study session notes from markdown strings to extract academic milestones."""
    milestones = []
    current_milestone: Dict[str, Any] = {}
    
    lines = content.splitlines()
    for line in lines:
        line_clean = line.strip().lstrip("-*+").strip()
        if not line_clean:
            continue
            
        subject_match = re.match(r"(?i)Subject\s*:\s*(.*)", line_clean)
        topic_match = re.match(r"(?i)Topic\s*:\s*(.*)", line_clean)
        status_match = re.match(r"(?i)Status\s*:\s*(.*)", line_clean)
        grade_match = re.match(r"(?i)Grade\s*:\s*(.*)", line_clean)
        
        if subject_match:
            if "subject" in current_milestone and "topic" in current_milestone:
                milestones.append(current_milestone)
                current_milestone = {}
            current_milestone["subject"] = subject_match.group(1).strip()
        elif topic_match:
            current_milestone["topic"] = topic_match.group(1).strip()
        elif status_match:
            current_milestone["status"] = status_match.group(1).strip()
        elif grade_match:
            try:
                current_milestone["grade"] = float(grade_match.group(1).strip())
            except ValueError:
                current_milestone["grade"] = 0.0
                
    if "subject" in current_milestone and "topic" in current_milestone:
        if "status" not in current_milestone:
            current_milestone["status"] = "Completed"
        if "grade" not in current_milestone:
            current_milestone["grade"] = 0.0
        milestones.append(current_milestone)
        
    return milestones

# Spaced Repetition SuperMemo SM-2 Interval Calculation Skill
def calculate_sm2_review(ease_factor: float, interval_days: int, consecutive_correct: int, quality: int) -> Tuple[float, int, str, int]:
    """
    Computes review tracking states based on the SuperMemo SM-2 Spaced Repetition Algorithm.
    """
    if quality < 3:
        new_consecutive_correct = 0
        new_interval_days = 1
        new_ease_factor = max(1.3, ease_factor - 0.2)
    else:
        new_consecutive_correct = consecutive_correct + 1
        
        if new_consecutive_correct == 1:
            new_interval_days = 1
        elif new_consecutive_correct == 2:
            new_interval_days = 6
        else:
            new_interval_days = int(round(interval_days * ease_factor))
            
        new_ease_factor = ease_factor + (0.1 - (5 - quality) * (0.08 + (5 - quality) * 0.02))
        new_ease_factor = max(1.3, new_ease_factor)
        
    today = datetime.date.today()
    next_date = today + datetime.timedelta(days=new_interval_days)
    
    return round(new_ease_factor, 2), new_interval_days, next_date.isoformat(), new_consecutive_correct


# ----------------- DOCUMENT PARSING & OCR SKILLS -----------------

def extract_text_from_pdf(file_path: str) -> str:
    """
    Extracts text from a PDF file.
    Tries native text extraction via pdfplumber first.
    If native text is empty (scanned image PDF), falls back to OCR via pytesseract if possible.
    """
    extracted_text = ""
    try:
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                text = page.extract_text()
                if text:
                    extracted_text += text + "\n"
    except Exception as e:
        print("pdfplumber extraction error:", e)
        
    # If text is too short or empty, try OCR (scanned document)
    if len(extracted_text.strip()) < 50:
        try:
            tesseract_path = r"C:\Program Files\Tesseract-OCR\tesseract.exe"
            if os.path.exists(tesseract_path):
                pytesseract.pytesseract.tesseract_cmd = tesseract_path
                
            # Convert PDF pages to images using pypdfium2 (since it was installed as pdfplumber dependency)
            import pypdfium2 as pdfium
            doc = pdfium.PdfDocument(file_path)
            ocr_text = ""
            for page in doc:
                bitmap = page.render(scale=2)
                pil_img = bitmap.to_pil()
                text = pytesseract.image_to_string(pil_img)
                if text:
                    ocr_text += text + "\n"
            if ocr_text.strip():
                extracted_text = "[OCR Extracted Text]:\n" + ocr_text
        except Exception as ocr_err:
            print("OCR extraction error:", ocr_err)
            
    return extracted_text

def extract_text_from_pptx(file_path: str) -> str:
    """
    Extracts text from a PPTX presentation file slide-by-slide.
    """
    extracted_text = ""
    try:
        prs = pptx.Presentation(file_path)
        for i, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            if slide_text:
                extracted_text += f"--- Slide {i+1} ---\n" + "\n".join(slide_text) + "\n\n"
    except Exception as e:
        print("pptx extraction error:", e)
    return extracted_text


# ----------------- 5-PHASE TABULAR MACHINE LEARNING PIPELINE -----------------

def inspect_data_integrity(df: pd.DataFrame) -> pd.DataFrame:
    """Performs deep diagnostic checks on the loaded dataset."""
    numeric_df = df.select_dtypes(include=[np.number])
    skewness = numeric_df.skew() if not numeric_df.empty else pd.Series(dtype='float64')
    
    diagnostics = pd.DataFrame({
        'Data Type': df.dtypes,
        'Missing Values': df.isnull().sum(),
        'Missing %': (df.isnull().sum() / len(df)) * 100,
        'Unique Values': df.nunique(),
        'Skewness': skewness
    })
    return diagnostics

def build_preprocessing_pipeline(num_cols: List[str], cat_cols: List[str]) -> ColumnTransformer:
    """Creates a reusable, leakage-free transformation pipeline."""
    num_transformer = Pipeline(steps=[
        ('imputer', SimpleImputer(strategy='median')),
        ('scaler', StandardScaler())
    ])
    
    cat_transformer = Pipeline(steps=[
        ('imputer', SimpleImputer(strategy='most_frequent')),
        ('onehot', OneHotEncoder(handle_unknown='ignore', sparse_output=False))
    ])
    
    preprocessor = ColumnTransformer(
        transformers=[
            ('num', num_transformer, num_cols),
            ('cat', cat_transformer, cat_cols)
        ]
    )
    return preprocessor

def evaluate_pipeline_cv(estimator: Any, preprocessor: ColumnTransformer, X: pd.DataFrame, y: pd.Series, cv_splits: int = 5) -> Dict[str, float]:
    """Runs a clean cross-validation loop safeguarding against leakage."""
    full_pipeline = Pipeline(steps=[
        ('preprocessor', preprocessor),
        ('model', estimator)
    ])
    
    cv = StratifiedKFold(n_splits=cv_splits, shuffle=True, random_state=42)
    scoring = ['accuracy', 'precision_weighted', 'recall_weighted', 'f1_weighted', 'roc_auc']
    
    scores = cross_validate(full_pipeline, X, y, cv=cv, scoring=scoring, n_jobs=-1, return_train_score=False)
    return {metric: float(np.mean(val)) for metric, val in scores.items()}

def get_baseline_models() -> Dict[str, Any]:
    """Returns initialized state-of-the-art tree architectures."""
    return {
        'XGBoost': XGBClassifier(n_estimators=100, learning_rate=0.05, max_depth=6, random_state=42, eval_metric='logloss'),
        'LightGBM': LGBMClassifier(n_estimators=100, learning_rate=0.05, max_depth=6, random_state=42, verbose=-1)
    }

def display_final_metrics(model_pipeline: Pipeline, X_test: pd.DataFrame, y_test: pd.Series, preprocessor: ColumnTransformer) -> Tuple[str, pd.DataFrame]:
    """Outputs granular metrics and uncovers structural feature importances."""
    preds = model_pipeline.predict(X_test)
    
    report = classification_report(y_test, preds)
    matrix = confusion_matrix(y_test, preds)
    
    output_text = f"📋 Classification Report:\n{report}\n\n📊 Confusion Matrix:\n{matrix}"
    
    feature_names = []
    try:
        feature_names.extend(preprocessor.transformers_[0][2])
        ohe = preprocessor.transformers_[1][1].named_steps['onehot']
        ohe_features = list(ohe.get_feature_names_out(preprocessor.transformers_[1][2]))
        feature_names.extend(ohe_features)
    except Exception:
        num_features = len(X_test.columns)
        feature_names = [f"Feature_{i}" for i in range(num_features)]
        
    model = model_pipeline.named_steps['model']
    raw_importances = model.feature_importances_
    
    if len(feature_names) < len(raw_importances):
        feature_names += [f"OHE_Feature_{i}" for i in range(len(raw_importances) - len(feature_names))]
    elif len(feature_names) > len(raw_importances):
        feature_names = feature_names[:len(raw_importances)]
        
    importance_df = pd.DataFrame({
        'Feature': feature_names,
        'Importance': raw_importances
    }).sort_values(by='Importance', ascending=False)
    
    return output_text, importance_df
